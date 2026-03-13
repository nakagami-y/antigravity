# -*- coding: utf-8 -*-
"""
PowerPoint \u751f\u6210\u30c6\u30f3\u30d7\u30ec\u30fc\u30c8 v2.0.0 (python-pptx + PIL)

\u4f7f\u3044\u65b9:
  1. SLIDES \u30ea\u30b9\u30c8\u306b\u30b9\u30e9\u30a4\u30c9\u5b9a\u7fa9\u3092\u8ffd\u52a0
  2. OUTPUT_PATH \u3092\u8a2d\u5b9a
  3. python template.py

13\u7a2e\u306e\u30b9\u30e9\u30a4\u30c9\u30bf\u30a4\u30d7:
  title, divider, fullbleed, agenda, content, two_column,
  three_cards, comparison, demo, table, triple_comparison,
  data, closing
"""

import os
import sys
import random
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData
except ImportError:
    print("pip install python-pptx")
    sys.exit(1)

try:
    from PIL import Image, ImageDraw, ImageFilter
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

# --- \u8a2d\u5b9a ---
OUTPUT_PATH = Path(__file__).parent.parent / "output" / "presentation.pptx"

# --- \u30b9\u30e9\u30a4\u30c9\u30b5\u30a4\u30ba ---
# \u30c7\u30d5\u30a9\u30eb\u30c8: 13.333 x 7.5 (16:9 \u30ef\u30a4\u30c9)
# \u4ee3\u66ff: 10.0 x 5.625 (16:9 \u53c2\u8003\u30c7\u30c3\u30ad\u5408\u308f\u305b)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
SW = SLIDE_WIDTH
SH = SLIDE_HEIGHT

# --- \u30ab\u30e9\u30fc\u30c6\u30fc\u30de ---
# \u30c7\u30d5\u30a9\u30eb\u30c8: Navy + Mint
C = {
    "bg": RGBColor(0xFF, 0xFF, 0xFF),
    "text": RGBColor(0x33, 0x33, 0x33),
    "navy": RGBColor(0x1E, 0x3A, 0x5F),
    "accent": RGBColor(0x00, 0xD4, 0xAA),
    "sub": RGBColor(0x66, 0x66, 0x66),
    "light_bg": RGBColor(0xF0, 0xF4, 0xF8),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "divider_sub": RGBColor(0xBB, 0xC3, 0xCF),
}

# \u4ee3\u66ff\u30c6\u30fc\u30de: Navy + Accent Blue (\u30bb\u30df\u30ca\u30fc\u5411\u3051)
# C = {
#     "bg": RGBColor(0x1B, 0x3A, 0x5C),
#     "text": RGBColor(0xFF, 0xFF, 0xFF),
#     "navy": RGBColor(0x1B, 0x3A, 0x5C),
#     "accent_blue": RGBColor(0x4A, 0x7F, 0xB5),
#     "white": RGBColor(0xFF, 0xFF, 0xFF),
#     "light_blue": RGBColor(0x8E, 0xAE, 0xD0),
#     "footer_blue": RGBColor(0x6B, 0x8D, 0xB5),
#     "page_num": RGBColor(0x4A, 0x6A, 0x8A),
# }

# --- \u30d5\u30a9\u30f3\u30c8 ---
# \u30c7\u30d5\u30a9\u30eb\u30c8: Meiryo
F_HEAD = "Meiryo"
F_BODY = "Meiryo"

# \u4ee3\u66ff: Georgia + Calibri (\u30bb\u30df\u30ca\u30fc\u5411\u3051)
# F_HEAD = "Georgia"
# F_BODY = "Calibri"


# --- \u30b9\u30e9\u30a4\u30c9\u5b9a\u7fa9\uff08\u3053\u3053\u3092\u7f6e\u304d\u63db\u3048\u308b\uff09 ---
SLIDES = [
    # --- title ---
    # {
    #     "type": "title",
    #     "title": "\u30d7\u30ec\u30bc\u30f3\u30bf\u30a4\u30c8\u30eb",
    #     "subtitle": "\u30b5\u30d6\u30bf\u30a4\u30c8\u30eb",
    #     "goal": "\u3053\u306e\u30bb\u30df\u30ca\u30fc\u306e\u30b4\u30fc\u30eb",  # optional
    #     "date": "2026-02-17",
    #     "image_path": "98_assets/mascot.png",  # optional: left-right split
    #     "note": "\u30b9\u30d4\u30fc\u30ab\u30fc\u30ce\u30fc\u30c8",
    # },
    # --- divider ---
    # {
    #     "type": "divider",
    #     "label": "Chapter 01",  # or chapter_number: 1
    #     "title": "\u7ae0\u30bf\u30a4\u30c8\u30eb",
    #     "subtitle": "\u3053\u306e\u7ae0\u306e\u6982\u8981",  # optional
    # },
    # --- fullbleed (\u30ad\u30fc\u30e1\u30c3\u30bb\u30fc\u30b8\u30fb\u3069\u771f\u3093\u4e2d) ---
    # {
    #     "type": "fullbleed",
    #     "main": "\u4e00\u6587\u30e1\u30c3\u30bb\u30fc\u30b8",
    #     "sub": "\u30b5\u30d6\u30c6\u30ad\u30b9\u30c8",  # optional
    #     "note": "\u30b9\u30d4\u30fc\u30ab\u30fc\u30ce\u30fc\u30c8",
    # },
    # --- agenda ---
    # {
    #     "type": "agenda",
    #     "title": "\u672c\u65e5\u306e\u30a2\u30b8\u30a7\u30f3\u30c0",
    #     "items": ["\u9805\u76ee1", "\u9805\u76ee2", "\u9805\u76ee3"],
    # },
    # --- content ---
    # {
    #     "type": "content",
    #     "title": "\u30b9\u30e9\u30a4\u30c9\u898b\u51fa\u3057",
    #     "bullets": ["\u30dd\u30a4\u30f3\u30c81", "\u30dd\u30a4\u30f3\u30c82", "\u30dd\u30a4\u30f3\u30c83"],
    #     "note": "\u30b9\u30d4\u30fc\u30ab\u30fc\u30ce\u30fc\u30c8",
    # },
    # --- two_column ---
    # {
    #     "type": "two_column",
    #     "title": "\u6bd4\u8f03",
    #     "left_title": "Before",
    #     "left_bullets": ["\u8ab2\u984c1", "\u8ab2\u984c2"],
    #     "right_title": "After",
    #     "right_bullets": ["\u89e3\u6c7a1", "\u89e3\u6c7a2"],
    #     "image_path": None,  # right side image
    # },
    # --- three_cards ---
    # {
    #     "type": "three_cards",
    #     "title": "3\u3064\u306e\u30dd\u30a4\u30f3\u30c8",
    #     "cards": [
    #         {"label": "01", "heading": "\u30ab\u30fc\u30c91", "body": "\u8aac\u660e\u6587"},
    #         {"label": "02", "heading": "\u30ab\u30fc\u30c92", "body": "\u8aac\u660e\u6587"},
    #         {"label": "03", "heading": "\u30ab\u30fc\u30c93", "body": "\u8aac\u660e\u6587"},
    #     ],
    # },
    # --- comparison (2\u629e\u6bd4\u8f03) ---
    # {
    #     "type": "comparison",
    #     "title": "\u6bd4\u8f03",
    #     "left_icon": "\u2717",
    #     "left_heading": "NG\u30d1\u30bf\u30fc\u30f3",
    #     "left_bullets": ["\u554f\u984c1", "\u554f\u984c2"],
    #     "right_icon": "\u25cb",
    #     "right_heading": "OK\u30d1\u30bf\u30fc\u30f3",
    #     "right_bullets": ["\u89e3\u6c7a1", "\u89e3\u6c7a2"],
    # },
    # --- demo ---
    # {
    #     "type": "demo",
    #     "title": "\u30c7\u30e2\u624b\u9806",
    #     "steps": [
    #         {"label": "Step 1", "text": "\u624b\u9806\u306e\u8aac\u660e"},
    #         {"label": "Step 2", "text": "\u624b\u9806\u306e\u8aac\u660e"},
    #     ],
    # },
    # --- table ---
    # {
    #     "type": "table",
    #     "title": "\u6a5f\u80fd\u6bd4\u8f03",
    #     "headers": ["\u9805\u76ee", "Claude Code", "\u5f93\u6765"],
    #     "rows": [
    #         ["\u9805\u76ee1", "\u25cb", "\u2717"],
    #         ["\u9805\u76ee2", "\u25cb", "\u25b3"],
    #     ],
    # },
    # --- triple_comparison (3\u5217\u6bd4\u8f03) ---
    # {
    #     "type": "triple_comparison",
    #     "title": "3\u3064\u306e\u9078\u629e\u80a2",
    #     "columns": [
    #         {"heading": "A", "bullets": ["\u7279\u5fb41", "\u7279\u5fb42"]},
    #         {"heading": "B", "bullets": ["\u7279\u5fb41", "\u7279\u5fb42"]},
    #         {"heading": "C", "bullets": ["\u7279\u5fb41", "\u7279\u5fb42"]},
    #     ],
    # },
    # --- data (\u8868 or \u30b0\u30e9\u30d5) ---
    # {
    #     "type": "data",
    #     "title": "\u58f2\u4e0a\u63a8\u79fb",
    #     "data_type": "bar",  # "table" | "bar" | "pie" | "line"
    #     "categories": ["1\u6708", "2\u6708", "3\u6708"],
    #     "series": [{"name": "\u58f2\u4e0a", "values": [100, 150, 200]}],
    # },
    # --- closing ---
    # {
    #     "type": "closing",
    #     "title": "\u3042\u308a\u304c\u3068\u3046\u3054\u3056\u3044\u307e\u3057\u305f",
    #     "message": "\u304a\u554f\u3044\u5408\u308f\u305b: info@example.com",
    #     "cta": "\u7121\u6599\u76f8\u8ac7\u306f\u3053\u3061\u3089",
    # },
]


# ===== \u30d8\u30eb\u30d1\u30fc =====

def _f(run, font_name, size, color, bold=False):
    """Set font on a run."""
    f = run.font
    f.name = font_name
    f.size = Pt(size)
    f.bold = bold
    if color:
        f.color.rgb = color


def _text(slide, left, top, width, height, text, font_name, size, color,
          bold=False, align=None, wrap=True, v_center=False):
    """Add a textbox with a single run. v_center=True sets vertical middle anchor."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = wrap
    if v_center:
        from pptx.oxml.ns import qn
        body_pr = tf._txBody.find(qn("a:bodyPr"))
        if body_pr is not None:
            body_pr.set("anchor", "ctr")
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    _f(r, font_name, size, color, bold)
    return tx


def _rect(slide, left, top, width, height, color):
    """Add a colored rectangle (accent line, bg, etc)."""
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _dark_bg(slide):
    """Navy full-slide background."""
    _rect(slide, Emu(0), Emu(0), SW, SH, C["navy"])


def _white_bg(slide):
    """White full-slide background."""
    _rect(slide, Emu(0), Emu(0), SW, SH, C["bg"])


def _note(slide, text):
    """Add speaker note."""
    if text:
        slide.notes_slide.notes_text_frame.text = text


def _page_num(slide, num, total):
    """Add page number (bottom-right). For large decks."""
    _text(slide, Inches(11.5), Inches(6.9), Inches(1.5), Inches(0.4),
          f"{num} / {total}", F_BODY, 9, C.get("page_num", C["sub"]),
          align=PP_ALIGN.RIGHT)


# ===== PIL \u753b\u50cf\u751f\u6210 =====

def _generate_title_bg(output_path, w=2000, h=1125):
    """Generate corporate-style title bg: skyline + diagonal geometric overlay.
    Returns path if successful, None if PIL unavailable."""
    if not HAS_PIL:
        return None
    random.seed(77)
    navy_dark = (15, 30, 52)
    navy_mid = (36, 74, 115)
    blue_accent = (74, 127, 181)

    base = Image.new("RGBA", (w, h), (*navy_dark, 255))
    draw = ImageDraw.Draw(base)

    # Left gradient for skyline area
    sky_w = int(w * 0.55)
    for y in range(h):
        ratio = 1.0 - y / h
        r = int(navy_dark[0] + 30 * ratio)
        g = int(navy_dark[1] + 50 * ratio)
        b = int(navy_dark[2] + 70 * ratio)
        draw.line([(0, y), (sky_w, y)], fill=(r, g, b, 255))

    # Abstract building silhouettes
    for _ in range(18):
        bx = random.randint(0, sky_w - 60)
        bw = random.randint(30, 80)
        bh = random.randint(120, 500)
        by = h - bh
        shade = random.randint(20, 50)
        draw.rectangle([bx, by, bx + bw, h], fill=(shade, shade + 15, shade + 35, 200))
        # Window lights
        for wy in range(by + 10, h - 10, 20):
            for wx in range(bx + 5, bx + bw - 5, 12):
                if random.random() > 0.6:
                    draw.rectangle([wx, wy, wx + 4, wy + 4],
                                   fill=(180, 200, 240, random.randint(80, 200)))

    # Diagonal polygon overlay (main separator)
    diag_pts = [(int(w * 0.42), 0), (w, 0), (w, h), (int(w * 0.58), h)]
    overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d2 = ImageDraw.Draw(overlay)
    d2.polygon(diag_pts, fill=(*navy_dark, 245))
    base = Image.alpha_composite(base, overlay)

    # Accent stripe
    stripe = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d3 = ImageDraw.Draw(stripe)
    sx = int(w * 0.40)
    d3.polygon([(sx, 0), (sx + 25, 0), (sx + 65, h), (sx + 40, h)],
               fill=(*blue_accent, 100))
    base = Image.alpha_composite(base, stripe)

    # Particles on right side
    particle = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    d4 = ImageDraw.Draw(particle)
    for _ in range(30):
        px = random.randint(int(w * 0.55), w - 20)
        py = random.randint(20, h - 20)
        pr = random.randint(1, 3)
        alpha = random.randint(40, 120)
        d4.ellipse([px - pr, py - pr, px + pr, py + pr],
                   fill=(*blue_accent, alpha))
    base = Image.alpha_composite(base, particle)

    # Bottom accent line
    final_draw = ImageDraw.Draw(base)
    final_draw.rectangle([0, h - 4, w, h], fill=(*blue_accent, 200))

    out = base.convert("RGB")
    out.save(str(output_path), "PNG", quality=95)
    return str(output_path)


def _generate_hero(output_path, w=2000, h=1125):
    """Generate abstract dark tech background with gradient, grid, network nodes."""
    if not HAS_PIL:
        return None
    random.seed(42)
    navy = (15, 30, 55)

    base = Image.new("RGBA", (w, h), (*navy, 255))
    draw = ImageDraw.Draw(base)

    # Gradient
    for y in range(h):
        ratio = y / h
        r = int(navy[0] * (1 - ratio * 0.3))
        g = int(navy[1] * (1 - ratio * 0.2))
        b = int(navy[2] + 30 * ratio)
        draw.line([(0, y), (w, y)], fill=(r, g, b, 255))

    # Subtle grid
    for x in range(0, w, 80):
        draw.line([(x, 0), (x, h)], fill=(40, 60, 90, 30), width=1)
    for y in range(0, h, 80):
        draw.line([(0, y), (w, y)], fill=(40, 60, 90, 30), width=1)

    # Network nodes with glow
    glow_layer = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow_layer)
    nodes = [(random.randint(50, w - 50), random.randint(50, h - 50)) for _ in range(15)]
    for nx, ny in nodes:
        gd.ellipse([nx - 15, ny - 15, nx + 15, ny + 15], fill=(74, 127, 181, 40))
        gd.ellipse([nx - 4, ny - 4, nx + 4, ny + 4], fill=(74, 127, 181, 180))
    glow_layer = glow_layer.filter(ImageFilter.GaussianBlur(radius=8))
    base = Image.alpha_composite(base, glow_layer)

    # Floating particles
    particle = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    pd = ImageDraw.Draw(particle)
    for _ in range(40):
        px = random.randint(0, w)
        py = random.randint(0, h)
        pr = random.randint(1, 3)
        pd.ellipse([px - pr, py - pr, px + pr, py + pr],
                   fill=(142, 174, 208, random.randint(30, 100)))
    base = Image.alpha_composite(base, particle)

    out = base.convert("RGB")
    out.save(str(output_path), "PNG", quality=95)
    return str(output_path)


# ===== \u30b9\u30e9\u30a4\u30c9\u30d3\u30eb\u30c0\u30fc =====

def build_title(prs, d):
    """type: title -- PIL bg or image split or text-only."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    img_path = d.get("image_path")
    has_image = img_path and Path(img_path).exists()
    use_pil_bg = d.get("pil_bg", False)

    if use_pil_bg and HAS_PIL:
        # Corporate-style: PIL background + text on right
        bg_path = OUTPUT_PATH.parent / "title_bg.png"
        bg = _generate_title_bg(bg_path)
        if bg:
            slide.shapes.add_picture(bg, Emu(0), Emu(0), SW, SH)
        else:
            _dark_bg(slide)
        # Text on right side of diagonal cut
        tx = Inches(5.20) if SW == Inches(10.0) else Inches(7.50)
        tw = Inches(4.40) if SW == Inches(10.0) else Inches(5.00)
        _text(slide, tx, Inches(1.00), tw, Inches(2.00),
              d.get("title", ""), F_HEAD, 40, C["white"], bold=True)
        _rect(slide, tx, Inches(3.10), Inches(2.50), Pt(3), C.get("accent_blue", C["accent"]))
        if d.get("subtitle"):
            _text(slide, tx, Inches(3.40), tw, Inches(0.80),
                  d["subtitle"], F_BODY, 16, C.get("light_blue", C["sub"]))
        if d.get("date"):
            _text(slide, tx, Inches(4.90), tw, Inches(0.30),
                  d["date"], F_BODY, 10, C.get("footer_blue", C["sub"]))
    elif has_image:
        # Left-right split: text left, image right
        _white_bg(slide)
        text_width = Inches(7)
        _text(slide, Inches(1), Inches(2), text_width, Inches(1.5),
              d.get("title", ""), F_HEAD, 36, C["navy"], bold=True)
        _rect(slide, Inches(1), Inches(3.6), Inches(3.5), Pt(4), C["accent"])
        y = Inches(3.9)
        if d.get("subtitle"):
            _text(slide, Inches(1), y, text_width, Inches(0.8),
                  d["subtitle"], F_BODY, 20, C["sub"])
            y += Inches(0.9)
        if d.get("date"):
            _text(slide, Inches(1), y, Inches(5), Inches(0.5),
                  d["date"], F_BODY, 14, C["sub"])
        slide.shapes.add_picture(str(img_path), Inches(8.5), Inches(1.5), width=Inches(4))
    else:
        # Text-only full-wide
        _white_bg(slide)
        tw = SW - Inches(2)
        _text(slide, Inches(1), Inches(2), tw, Inches(1.5),
              d.get("title", ""), F_HEAD, 36, C["navy"], bold=True)
        _rect(slide, Inches(1), Inches(3.6), Inches(3.5), Pt(4), C["accent"])
        y = Inches(3.9)
        if d.get("subtitle"):
            _text(slide, Inches(1), y, tw, Inches(0.8),
                  d["subtitle"], F_BODY, 20, C["sub"])
            y += Inches(0.9)
        if d.get("goal"):
            tx = _text(slide, Inches(1), y, tw, Inches(0.6),
                       "", F_BODY, 16, C["navy"])
            tf = tx.text_frame
            p = tf.paragraphs[0]
            r1 = p.add_run()
            r1.text = "\u25b8 Goal: "
            _f(r1, F_BODY, 16, C["accent"], bold=True)
            r2 = p.add_run()
            r2.text = d["goal"]
            _f(r2, F_BODY, 16, C["navy"])
            y += Inches(0.7)
        if d.get("date"):
            _text(slide, Inches(1), y, Inches(5), Inches(0.5),
                  d["date"], F_BODY, 14, C["sub"])
        # Bottom accent line
        _rect(slide, Inches(1), Inches(6.8), tw, Pt(4), C["accent"])

    _note(slide, d.get("note", ""))
    return slide


def build_divider(prs, d):
    """type: divider -- Navy bg, centered label + title + accent line."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    slide_h = SH.inches

    label = d.get("label", "")
    if not label and d.get("chapter_number"):
        label = f"Chapter {d['chapter_number']:02d}"

    # Centered layout
    cw = SW - Inches(2)  # content width
    if label:
        _text(slide, Inches(1), Inches(slide_h * 0.33), cw, Inches(0.5),
              label, F_BODY, 14, C.get("accent_blue", C["accent"]),
              align=PP_ALIGN.CENTER)
    _text(slide, Inches(1), Inches(slide_h * 0.40), cw, Inches(1.2),
          d.get("title", ""), F_HEAD, 30, C["white"], bold=True,
          align=PP_ALIGN.CENTER)
    # Accent line (centered)
    line_w = Inches(2.40)
    line_left = Inches((SW.inches - 2.40) / 2)
    _rect(slide, line_left, Inches(slide_h * 0.58), line_w, Pt(3),
          C.get("accent_blue", C["accent"]))
    if d.get("subtitle"):
        _text(slide, Inches(1), Inches(slide_h * 0.65), cw, Inches(0.8),
              d["subtitle"], F_BODY, 18, C["divider_sub"],
              align=PP_ALIGN.CENTER)
    _note(slide, d.get("note", ""))
    return slide


def build_fullbleed(prs, d):
    """type: fullbleed -- Dark bg with dead-center text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _dark_bg(slide)
    slide_h = SH.inches

    main = d.get("main", "")
    has_sub = bool(d.get("sub"))
    cw = SW - Inches(1.2)

    if has_sub:
        total = 2.80
        block_start = (slide_h - total) / 2
        _text(slide, Inches(0.60), Inches(block_start), cw, Inches(2.00),
              main, F_HEAD, 30, C["white"], bold=True,
              align=PP_ALIGN.CENTER, v_center=True)
        sub_y = block_start + 2.00 + 0.20
        _text(slide, Inches(0.60), Inches(sub_y), cw, Inches(0.60),
              d["sub"], F_BODY, 14, C.get("light_blue", C["divider_sub"]),
              align=PP_ALIGN.CENTER, v_center=True)
    else:
        total = 2.00
        block_start = (slide_h - total) / 2
        _text(slide, Inches(0.60), Inches(block_start), cw, Inches(2.00),
              main, F_HEAD, 30, C["white"], bold=True,
              align=PP_ALIGN.CENTER, v_center=True)

    _note(slide, d.get("note", ""))
    return slide


def build_agenda(prs, d):
    """type: agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _text(slide, Inches(1), Inches(0.6), SW - Inches(2), Inches(0.8),
          d.get("title", "\u30a2\u30b8\u30a7\u30f3\u30c0"), F_HEAD, 28, C["navy"], bold=True)
    _rect(slide, Inches(1), Inches(1.5), Inches(3), Pt(4), C["accent"])

    items = d.get("items", [])
    txbox = slide.shapes.add_textbox(Inches(1.5), Inches(2), SW - Inches(3), Inches(4.5))
    tf = txbox.text_frame
    tf.word_wrap = True
    spacing = min(Pt(16), Pt(max(8, int(72 / max(len(items), 1)))))
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        r1 = p.add_run()
        r1.text = f"{i + 1}.  "
        _f(r1, F_BODY, 20, C["accent"], bold=True)
        r2 = p.add_run()
        r2.text = item
        _f(r2, F_BODY, 20, C["text"])
    _note(slide, d.get("note", ""))
    return slide


def build_content(prs, d):
    """type: content"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _text(slide, Inches(1), Inches(0.6), SW - Inches(2), Inches(0.8),
          d.get("title", ""), F_HEAD, 28, C["navy"], bold=True)
    _rect(slide, Inches(1), Inches(1.5), Inches(3), Pt(4), C["accent"])

    bullets = d.get("bullets", [])
    txbox = slide.shapes.add_textbox(Inches(1.2), Inches(2), SW - Inches(2.4), Inches(4.5))
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(12)
        r1 = p.add_run()
        r1.text = "\u25b8  "
        _f(r1, F_BODY, 16, C["accent"])
        r2 = p.add_run()
        r2.text = bullet
        _f(r2, F_BODY, 16, C["text"])
    _note(slide, d.get("note", ""))
    return slide


def build_two_column(prs, d):
    """type: two_column"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _text(slide, Inches(1), Inches(0.6), SW - Inches(2), Inches(0.8),
          d.get("title", ""), F_HEAD, 28, C["navy"], bold=True)
    _rect(slide, Inches(1), Inches(1.5), Inches(3), Pt(4), C["accent"])

    col_w = (SW.inches - 3) / 2
    left_x = Inches(1)
    right_x = Inches(1 + col_w + 1)

    # Left column
    if d.get("left_title"):
        _text(slide, left_x, Inches(2), Inches(col_w), Inches(0.5),
              d["left_title"], F_HEAD, 20, C["navy"], bold=True)
    if d.get("left_bullets"):
        txbox = slide.shapes.add_textbox(left_x, Inches(2.7), Inches(col_w), Inches(4))
        tf = txbox.text_frame
        tf.word_wrap = True
        for i, b in enumerate(d["left_bullets"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            r1 = p.add_run()
            r1.text = "\u25b8  "
            _f(r1, F_BODY, 16, C["accent"])
            r2 = p.add_run()
            r2.text = b
            _f(r2, F_BODY, 16, C["text"])

    # Right column
    img_path = d.get("image_path")
    if img_path and Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), right_x, Inches(2),
                                 width=Inches(col_w))
    else:
        if d.get("right_title"):
            _text(slide, right_x, Inches(2), Inches(col_w), Inches(0.5),
                  d["right_title"], F_HEAD, 20, C.get("accent_blue", C["accent"]), bold=True)
        if d.get("right_bullets"):
            txbox = slide.shapes.add_textbox(right_x, Inches(2.7), Inches(col_w), Inches(4))
            tf = txbox.text_frame
            tf.word_wrap = True
            for i, b in enumerate(d["right_bullets"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(10)
                r1 = p.add_run()
                r1.text = "\u25b8  "
                _f(r1, F_BODY, 16, C["accent"])
                r2 = p.add_run()
                r2.text = b
                _f(r2, F_BODY, 16, C["text"])
    _note(slide, d.get("note", ""))
    return slide


def build_three_cards(prs, d):
    """type: three_cards -- 3 column cards with accent bar."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _text(slide, Inches(1), Inches(0.6), SW - Inches(2), Inches(0.8),
          d.get("title", ""), F_HEAD, 28, C["navy"], bold=True)

    cards = d.get("cards", [])
    n = len(cards)
    if n == 0:
        return slide
    margin = 1.0
    gap = 0.30
    total_w = SW.inches - margin * 2
    card_w = (total_w - gap * (n - 1)) / n
    card_h = 3.80

    for i, card in enumerate(cards):
        cx = Inches(margin + i * (card_w + gap))
        cy = Inches(2.00)
        # Card bg
        _rect(slide, cx, cy, Inches(card_w), Inches(card_h), C["light_bg"])
        # Top accent bar
        _rect(slide, cx, cy, Inches(card_w), Pt(4), C["accent"])
        # Label
        if card.get("label"):
            _text(slide, cx + Inches(0.30), cy + Inches(0.30),
                  Inches(card_w - 0.60), Inches(0.40),
                  card["label"], F_BODY, 14, C["accent"], bold=True)
        # Heading
        if card.get("heading"):
            _text(slide, cx + Inches(0.30), cy + Inches(0.80),
                  Inches(card_w - 0.60), Inches(0.50),
                  card["heading"], F_HEAD, 18, C["navy"], bold=True)
        # Body
        if card.get("body"):
            _text(slide, cx + Inches(0.30), cy + Inches(1.40),
                  Inches(card_w - 0.60), Inches(card_h - 1.80),
                  card["body"], F_BODY, 13, C["text"])
    _note(slide, d.get("note", ""))
    return slide


def build_comparison(prs, d):
    """type: comparison -- 2 side-by-side cards (X vs O)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _text(slide, Inches(1), Inches(0.6), SW - Inches(2), Inches(0.8),
          d.get("title", ""), F_HEAD, 28, C["navy"], bold=True)

    margin = 1.0
    gap = 0.50
    total_w = SW.inches - margin * 2
    card_w = (total_w - gap) / 2
    card_h = 4.20

    for i, side in enumerate(["left", "right"]):
        cx = Inches(margin + i * (card_w + gap))
        cy = Inches(2.00)
        bg_color = C["light_bg"]
        _rect(slide, cx, cy, Inches(card_w), Inches(card_h), bg_color)
        # Icon
        icon = d.get(f"{side}_icon", "")
        if icon:
            icon_color = RGBColor(0xCC, 0x33, 0x33) if i == 0 else C["accent"]
            _text(slide, cx + Inches(0.30), cy + Inches(0.20),
                  Inches(0.60), Inches(0.60),
                  icon, F_HEAD, 28, icon_color, bold=True)
        # Heading
        heading = d.get(f"{side}_heading", "")
        if heading:
            _text(slide, cx + Inches(1.00), cy + Inches(0.25),
                  Inches(card_w - 1.40), Inches(0.50),
                  heading, F_HEAD, 20, C["navy"], bold=True)
        # Bullets
        bullets = d.get(f"{side}_bullets", [])
        if bullets:
            txbox = slide.shapes.add_textbox(
                cx + Inches(0.50), cy + Inches(1.00),
                Inches(card_w - 1.00), Inches(card_h - 1.40))
            tf = txbox.text_frame
            tf.word_wrap = True
            for j, b in enumerate(bullets):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_after = Pt(8)
                r = p.add_run()
                r.text = f"\u25b8 {b}"
                _f(r, F_BODY, 14, C["text"])
    _note(slide, d.get("note", ""))
    return slide


def build_demo(prs, d):
    """type: demo -- Numbered steps."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _text(slide, Inches(1), Inches(0.6), SW - Inches(2), Inches(0.8),
          d.get("title", ""), F_HEAD, 28, C["navy"], bold=True)
    _rect(slide, Inches(1), Inches(1.5), Inches(3), Pt(4), C["accent"])

    steps = d.get("steps", [])
    n = len(steps)
    available = 4.50
    spacing = min(0.90, available / max(n, 1))
    y = 2.20
    for step in steps:
        label = step.get("label", "")
        text = step.get("text", "")
        # Step label (accent bg circle-ish)
        _text(slide, Inches(1.2), Inches(y), Inches(1.2), Inches(0.40),
              label, F_BODY, 14, C["accent"], bold=True)
        # Step text
        _text(slide, Inches(2.60), Inches(y), SW - Inches(3.80), Inches(0.40),
              text, F_BODY, 15, C["text"])
        y += spacing
    _note(slide, d.get("note", ""))
    return slide


def build_table(prs, d):
    """type: table -- Styled data table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _text(slide, Inches(1), Inches(0.6), SW - Inches(2), Inches(0.8),
          d.get("title", ""), F_HEAD, 28, C["navy"], bold=True)

    headers = d.get("headers", [])
    rows = d.get("rows", [])
    n_rows = len(rows) + 1
    n_cols = len(headers) if headers else (len(rows[0]) if rows else 1)

    tbl_shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(1), Inches(2), SW - Inches(2), Inches(4.5))
    table = tbl_shape.table

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                _f(r, F_BODY, 14, C["white"], bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["navy"]

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    _f(r, F_BODY, 13, C["text"])
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["light_bg"]
    _note(slide, d.get("note", ""))
    return slide


def build_triple_comparison(prs, d):
    """type: triple_comparison -- 3 column comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _text(slide, Inches(1), Inches(0.6), SW - Inches(2), Inches(0.8),
          d.get("title", ""), F_HEAD, 28, C["navy"], bold=True)

    columns = d.get("columns", [])
    n = len(columns)
    if n == 0:
        return slide
    margin = 1.0
    gap = 0.30
    total_w = SW.inches - margin * 2
    col_w = (total_w - gap * (n - 1)) / n
    col_h = 4.00

    for i, col in enumerate(columns):
        cx = Inches(margin + i * (col_w + gap))
        cy = Inches(2.20)
        _rect(slide, cx, cy, Inches(col_w), Inches(col_h), C["light_bg"])
        _rect(slide, cx, cy, Inches(col_w), Pt(4), C["accent"])
        if col.get("heading"):
            _text(slide, cx + Inches(0.20), cy + Inches(0.20),
                  Inches(col_w - 0.40), Inches(0.50),
                  col["heading"], F_HEAD, 18, C["navy"], bold=True,
                  align=PP_ALIGN.CENTER)
        if col.get("bullets"):
            txbox = slide.shapes.add_textbox(
                cx + Inches(0.30), cy + Inches(0.90),
                Inches(col_w - 0.60), Inches(col_h - 1.20))
            tf = txbox.text_frame
            tf.word_wrap = True
            for j, b in enumerate(col["bullets"]):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.space_after = Pt(8)
                r = p.add_run()
                r.text = f"\u25b8 {b}"
                _f(r, F_BODY, 13, C["text"])
    _note(slide, d.get("note", ""))
    return slide


def build_data(prs, d):
    """type: data -- chart (bar/pie/line) or table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _text(slide, Inches(1), Inches(0.6), SW - Inches(2), Inches(0.8),
          d.get("title", ""), F_HEAD, 28, C["navy"], bold=True)
    _rect(slide, Inches(1), Inches(1.5), Inches(3), Pt(4), C["accent"])

    data_type = d.get("data_type", "table")
    if data_type == "table":
        # Inline table
        headers = d.get("headers", [])
        rows = d.get("rows", [])
        if headers and rows:
            n_rows = len(rows) + 1
            n_cols = len(headers)
            tbl = slide.shapes.add_table(
                n_rows, n_cols, Inches(1), Inches(2), SW - Inches(2), Inches(4.5))
            table = tbl.table
            for ci, h in enumerate(headers):
                cell = table.cell(0, ci)
                cell.text = h
                cell.fill.solid()
                cell.fill.fore_color.rgb = C["navy"]
            for ri, row in enumerate(rows):
                for ci, val in enumerate(row):
                    table.cell(ri + 1, ci).text = str(val)
    else:
        chart_types = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "pie": XL_CHART_TYPE.PIE,
            "line": XL_CHART_TYPE.LINE,
        }
        xl_type = chart_types.get(data_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
        chart_data = CategoryChartData()
        chart_data.categories = d.get("categories", [])
        for s in d.get("series", []):
            chart_data.add_series(s["name"], s["values"])
        slide.shapes.add_chart(
            xl_type, Inches(1.5), Inches(2), SW - Inches(3), Inches(4.8), chart_data)
    _note(slide, d.get("note", ""))
    return slide


def build_closing(prs, d):
    """type: closing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _white_bg(slide)
    _text(slide, Inches(1), Inches(2.5), SW - Inches(2), Inches(1.2),
          d.get("title", "\u3042\u308a\u304c\u3068\u3046\u3054\u3056\u3044\u307e\u3057\u305f"),
          F_HEAD, 36, C["navy"], bold=True, align=PP_ALIGN.CENTER)
    _rect(slide, Inches((SW.inches - 4) / 2), Inches(3.8), Inches(4), Pt(4), C["accent"])
    if d.get("message"):
        _text(slide, Inches(1), Inches(4.2), SW - Inches(2), Inches(0.8),
              d["message"], F_BODY, 18, C["sub"], align=PP_ALIGN.CENTER)
    if d.get("cta"):
        btn_w = Inches(4.333)
        btn_left = Inches((SW.inches - 4.333) / 2)
        shape = slide.shapes.add_shape(5, btn_left, Inches(5.3), btn_w, Inches(0.7))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C["accent"]
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = d["cta"]
        _f(r, F_BODY, 18, C["white"], bold=True)
    _note(slide, d.get("note", ""))
    return slide


# ===== \u30d3\u30eb\u30c0\u30fc\u30c7\u30a3\u30b9\u30d1\u30c3\u30c1 =====

BUILDERS = {
    "title": build_title,
    "divider": build_divider,
    "fullbleed": build_fullbleed,
    "agenda": build_agenda,
    "content": build_content,
    "two_column": build_two_column,
    "three_cards": build_three_cards,
    "comparison": build_comparison,
    "demo": build_demo,
    "table": build_table,
    "triple_comparison": build_triple_comparison,
    "data": build_data,
    "closing": build_closing,
}


def main():
    if not SLIDES:
        print("No slides defined. Edit the SLIDES list.")
        sys.exit(0)

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    total = len(SLIDES)
    print(f"slides: {total}")

    for i, slide_data in enumerate(SLIDES):
        stype = slide_data.get("type", "content")
        builder = BUILDERS.get(stype)
        if builder:
            builder(prs, slide_data)
            title = slide_data.get("title", slide_data.get("main", "(no title)"))
            print(f"  [{i + 1}] {stype}: {title}")
        else:
            print(f"  [{i + 1}] unknown type: {stype} (skipped)")

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    size_kb = OUTPUT_PATH.stat().st_size / 1024
    print(f"\nsaved: {OUTPUT_PATH}")
    print(f"size: {size_kb:.1f} KB")


if __name__ == "__main__":
    main()
